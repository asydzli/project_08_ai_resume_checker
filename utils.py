from pypdf import PdfReader
from docx import Document
from pptx import Presentation


def read_pdf(file):

    reader = PdfReader(file)

    text = ""

    for page in reader.pages:
        page_text = page.extract_text()

        if page_text:
            text += page_text + "\n"

    return text


def read_docx(file):

    doc = Document(file)

    text = ""

    for para in doc.paragraphs:
        text += para.text + "\n"

    return text


def read_pptx(file):

    prs = Presentation(file)

    text = ""

    for slide in prs.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                text += shape.text + "\n"

    return text